"""
Genera Evaluacion_1_Presentacion.pptx v4 — Con fórmulas LaTeX renderizadas como imágenes,
analogías en Teledetección, impacto tecnológico, y contenido mejorado.
"""
import os, sys, io, tempfile
sys.stdout.reconfigure(encoding='utf-8')
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Rutas ──
SCRATCH   = r"C:\Users\jorge\.gemini\antigravity\brain\91c9ba88-7aab-4773-945b-d9e180a70db0\scratch"
OUT_PPTX  = r"D:\00_FisicaModerna\02_TeoriaCuanticaTemprana\Evaluacion\Evaluacion_1_Presentacion.pptx"
CHART_RJ  = os.path.join(SCRATCH, "chart_rj_planck.png")
CHART_BB  = os.path.join(SCRATCH, "chart_blackbody.png")
FORMULAS  = os.path.join(SCRATCH, "formulas")
os.makedirs(FORMULAS, exist_ok=True)

# ── Paleta ──
BG      = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT1 = RGBColor(0x00, 0xC8, 0xFF)   # cyan
ACCENT2 = RGBColor(0xFF, 0x6B, 0x35)   # orange
ACCENT3 = RGBColor(0x7B, 0xFF, 0xC4)   # green
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xAA, 0xBB)
YELLOW  = RGBColor(0xFF, 0xE0, 0x66)
PURPLE  = RGBColor(0xC0, 0x84, 0xFC)
DARKBG  = RGBColor(0x12, 0x20, 0x33)

# ── LaTeX Formula Renderer ──
def render_latex(latex_str, filename, fontsize=18, color='#FFE066', dpi=200):
    """Render a LaTeX formula to a PNG image with transparent background."""
    path = os.path.join(FORMULAS, filename)
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.patch.set_alpha(0)
    text = fig.text(0, 0.5, f'${latex_str}$', fontsize=fontsize, color=color,
                    ha='left', va='center',
                    usetex=False,  # Use mathtext, not full LaTeX
                    fontfamily='serif')
    fig.savefig(path, dpi=dpi, bbox_inches='tight', pad_inches=0.06,
                transparent=True, facecolor='none')
    plt.close(fig)
    return path

# Pre-render all formulas
def render_all_formulas():
    formulas = {
        # Slide 2
        'stefan':     (r'R = \sigma\, T^4', 20, '#FFE066'),
        'sigma_val':  (r'\sigma = 5.670 \times 10^{-8}\;\mathrm{W\,m^{-2}\,K^{-4}}', 14, '#AAAABB'),
        'wien':       (r'\lambda_{\max} \cdot T = 2.898 \times 10^{-3}\;\mathrm{m \cdot K}', 18, '#FFE066'),
        # Slide 3
        'equip':      (r'\langle E \rangle = k_B T \quad \mathrm{por\;modo}', 18, '#FFFFFF'),
        'modes':      (r'g(\nu) = \dfrac{8\pi\nu^2}{c^3}', 20, '#FFFFFF'),
        'rj':         (r'u(\nu, T) = \dfrac{8\pi\nu^2}{c^3}\, k_B T', 22, '#FFE066'),
        'diverge':    (r'\int_0^\infty u(\nu,T)\,d\nu \;\to\; \infty', 20, '#FF6B35'),
        # Slide 4
        'quant':      (r'E_n = n\,h\nu \quad (n = 0, 1, 2, 3, \ldots)', 22, '#FFE066'),
        'h_val':      (r'h = 6.626 \times 10^{-34}\;\mathrm{J \cdot s}', 14, '#AAAABB'),
        'planck':     (r'u(\nu, T) = \dfrac{8\pi h\nu^3}{c^3}\cdot\dfrac{1}{e^{h\nu/k_BT}-1}', 22, '#FFE066'),
        'mean_e':     (r'\bar{E} = \dfrac{h\nu}{e^{h\nu/k_BT}-1}', 20, '#FFE066'),
        'low_freq':   (r'h\nu \ll k_BT \;\Rightarrow\; \bar{E} \approx k_BT \;\checkmark', 16, '#7BFFC4'),
        'high_freq':  (r'h\nu \gg k_BT \;\Rightarrow\; \bar{E} \to 0 \;\checkmark', 16, '#7BFFC4'),
    }
    paths = {}
    for key, (latex, fs, col) in formulas.items():
        paths[key] = render_latex(latex, f'{key}.png', fontsize=fs, color=col)
        print(f'  Formula: {key}')
    return paths

# ── Helpers ──
def set_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG

def add_text(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    p = tf.add_paragraph(); p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return p

def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def formula_img(slide, key, left, top, width, height):
    """Insert a pre-rendered formula image."""
    return slide.shapes.add_picture(F[key], Inches(left), Inches(top), Inches(width), Inches(height))

def header(slide, text, color=ACCENT1):
    rect(slide, 0, 0, 13.33, 0.08, color)
    tb_tag = textbox(slide, 0.4, 0.15, 12.0, 0.35)
    add_text(tb_tag.text_frame, "DIPLOMADO EN FÍSICA MODERNA — EVALUACIÓN 1", 9, color=GRAY, bold=True)
    tb = textbox(slide, 0.4, 0.42, 12.0, 0.7)
    add_text(tb.text_frame, text, 24, bold=True, color=color)
    rect(slide, 0.4, 1.05, 12.5, 0.03, color)
    rect(slide, 0, 7.32, 13.33, 0.18, color)

# ── Gráficas ──
def make_blackbody_chart():
    h, c, kB = 6.626e-34, 3e8, 1.381e-23
    lam = np.linspace(100e-9, 3000e-9, 2000)
    fig, ax = plt.subplots(figsize=(6.5, 3.8), facecolor='#0D1B2A')
    ax.set_facecolor('#0D1B2A')
    for T, col, lbl in [(3000,'#FF6B35','3 000 K'), (4500,'#FFE066','4 500 K'), (6000,'#00C8FF','6 000 K')]:
        u = (8*np.pi*h*c / lam**5) / (np.exp(h*c/(lam*kB*T)) - 1)
        u /= u.max()
        ax.plot(lam*1e9, u, color=col, lw=2.4, label=lbl)
        lmax = lam[np.argmax(u)]*1e9
        ax.axvline(lmax, color=col, lw=0.8, ls='--', alpha=0.4)
        ax.annotate(f'{lmax:.0f} nm', xy=(lmax, 1.0), xytext=(lmax+80, 0.92),
                    fontsize=7, color=col, alpha=0.8)
    ax.axvspan(380, 700, alpha=0.10, color='white')
    ax.text(540, 0.97, 'Visible', color='white', fontsize=8, ha='center', alpha=0.6,
            bbox=dict(boxstyle='round,pad=0.2', facecolor='#0D1B2A', edgecolor='#334455', alpha=0.8))
    ax.set_xlabel('Longitud de onda λ (nm)', color='#AAAABB', fontsize=10)
    ax.set_ylabel('Intensidad relativa', color='#AAAABB', fontsize=10)
    ax.set_title('Espectro del Cuerpo Negro', color='white', fontsize=12, pad=10, fontweight='bold')
    ax.tick_params(colors='#AAAABB', labelsize=8)
    for s in ax.spines.values(): s.set_edgecolor('#334455')
    ax.legend(fontsize=8.5, facecolor='#0D1B2A', edgecolor='#445566',
              labelcolor='white', loc='upper right', framealpha=0.9)
    ax.set_xlim(100, 3000); ax.set_ylim(0, 1.08)
    plt.tight_layout()
    plt.savefig(CHART_BB, dpi=200, bbox_inches='tight', facecolor='#0D1B2A')
    plt.close()
    print(f"  Chart BB -> {CHART_BB}")

def make_rj_planck_chart():
    h, c, kB, T = 6.626e-34, 3e8, 1.381e-23, 5000
    nu = np.linspace(1e12, 3e14, 3000)
    planck = (8*np.pi*h*nu**3/c**3) / (np.exp(h*nu/(kB*T))-1)
    rj = (8*np.pi*nu**2/c**3) * kB*T
    norm = planck.max()
    planck /= norm; rj /= norm; rj = np.clip(rj, 0, 4.5)
    fig, ax = plt.subplots(figsize=(6.5, 3.8), facecolor='#0D1B2A')
    ax.set_facecolor('#0D1B2A')
    ax.fill_between(nu/1e13, planck, alpha=0.15, color='#00C8FF')
    ax.plot(nu/1e13, planck, color='#00C8FF', lw=2.8, label='Planck (experimental)', zorder=3)
    ax.plot(nu/1e13, rj, color='#FF6B35', lw=2.2, ls='--', label='Rayleigh-Jeans (clásica)', zorder=2)
    mask = rj > 1.15
    ax.fill_between(nu[mask]/1e13, rj[mask], planck[mask], color='#FF6B35', alpha=0.12)
    mid_idx = len(nu[mask])//2
    ax.annotate('CATÁSTROFE UV\n(divergencia → ∞)',
                xy=(nu[mask][mid_idx]/1e13, np.clip(rj[mask][mid_idx], 0, 3.5)),
                fontsize=9, color='#FF6B35', ha='center', fontweight='bold',
                arrowprops=dict(arrowstyle='->', color='#FF6B35', lw=1.5),
                xytext=(nu[mask][mid_idx]/1e13 - 2, 3.8))
    # Zona de acuerdo
    ax.annotate('Zona de acuerdo\n(baja frecuencia)',
                xy=(3, 0.25), fontsize=7.5, color='#7BFFC4', ha='center',
                fontstyle='italic', alpha=0.8)
    ax.set_xlabel('Frecuencia ν (×10¹³ Hz)', color='#AAAABB', fontsize=10)
    ax.set_ylabel('Densidad de energía u(ν, T)', color='#AAAABB', fontsize=10)
    ax.set_title('Rayleigh-Jeans vs. Planck  (T = 5 000 K)', color='white', fontsize=12, pad=10, fontweight='bold')
    ax.tick_params(colors='#AAAABB', labelsize=8)
    for s in ax.spines.values(): s.set_edgecolor('#334455')
    ax.legend(fontsize=8.5, facecolor='#0D1B2A', edgecolor='#445566',
              labelcolor='white', loc='upper left', framealpha=0.9)
    ax.set_xlim(0, 30); ax.set_ylim(0, 4.3)
    plt.tight_layout()
    plt.savefig(CHART_RJ, dpi=200, bbox_inches='tight', facecolor='#0D1B2A')
    plt.close()
    print(f"  Chart RJ -> {CHART_RJ}")

# ══════════════════════════════════════════════════════════════════════════
def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ═══════ SLIDE 1: PORTADA ═══════
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    rect(s1, 0, 0, 13.33, 0.08, ACCENT1)

    s1.shapes.add_picture(CHART_BB, Inches(6.6), Inches(0.8), Inches(6.5), Inches(4.0))

    tb = textbox(s1, 0.5, 0.9, 6.2, 1.2)
    tf = tb.text_frame; tf.word_wrap = True
    add_text(tf, "El Cuerpo Negro", 40, bold=True, color=ACCENT1)

    tb2 = textbox(s1, 0.5, 2.1, 6.2, 0.8)
    tf2 = tb2.text_frame; tf2.word_wrap = True
    add_text(tf2, "y la Revolución de Planck", 30, bold=True, color=WHITE)

    tb3 = textbox(s1, 0.5, 3.1, 6.2, 0.7)
    tf3 = tb3.text_frame; tf3.word_wrap = True
    add_text(tf3, "De la crisis del paradigma clásico\nal nacimiento de la Física Cuántica", 15, color=GRAY, italic=True)

    rect(s1, 0.5, 4.0, 5.0, 0.03, ACCENT1)

    # Info box
    rect(s1, 0.5, 4.2, 5.8, 2.4, DARKBG)
    tb4 = textbox(s1, 0.7, 4.3, 5.4, 2.2)
    tf4 = tb4.text_frame; tf4.word_wrap = True
    add_text(tf4, "Formato", 10, bold=True, color=ACCENT1)
    add_text(tf4, "Exposición oral en video (5 minutos máx.)", 11, color=WHITE)
    add_text(tf4, "", 3)
    add_text(tf4, "Enfoque pedagógico", 10, bold=True, color=ACCENT1)
    add_text(tf4, "Transposición didáctica con analogías en", 11, color=WHITE)
    add_text(tf4, "Geomática y Teledetección", 11, bold=True, color=ACCENT3)
    add_text(tf4, "", 3)
    add_text(tf4, "Tópicos obligatorios", 10, bold=True, color=ACCENT1)
    add_text(tf4, "1. Radiación Térmica del Cuerpo Negro", 11, color=WHITE)
    add_text(tf4, "2. Catástrofe Ultravioleta", 11, color=WHITE)
    add_text(tf4, "3. Cuantización de Planck", 11, color=WHITE)
    add_text(tf4, "", 3)
    add_text(tf4, "Evaluación 1 — Teoría Cuántica Temprana", 10, color=GRAY)
    add_text(tf4, "Docentes: Pablo Solano · Paulraj Manidurai · 2026", 10, color=GRAY)

    rect(s1, 0, 7.32, 13.33, 0.18, ACCENT1)

    # ═══════ SLIDE 2: CUERPO NEGRO ═══════
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    header(s2, "1. La Radiación del Cuerpo Negro e Incapacidad Clásica", ACCENT1)

    s2.shapes.add_picture(CHART_BB, Inches(6.4), Inches(1.15), Inches(6.6), Inches(3.8))

    tb2 = textbox(s2, 0.4, 1.15, 6.0, 1.6)
    tf2 = tb2.text_frame; tf2.word_wrap = True
    add_text(tf2, "Concepto y Propiedades Físicas", 14, bold=True, color=ACCENT3)
    add_text(tf2, "• Cuerpo Negro Ideal: Absorbe el 100% de la", 12, color=WHITE)
    add_text(tf2, "  radiación incidente y emite solo según T  (ε = 1).", 12, color=WHITE)
    add_text(tf2, "• Modelo Experimental: Cavidad cerrada con", 12, color=WHITE)
    add_text(tf2, "  un pequeño orificio de salida.", 12, color=WHITE)

    # Formulas rendered as images
    add_text(tf2, "", 4)
    tb_laws = textbox(s2, 0.4, 2.9, 6.0, 0.3)
    add_text(tb_laws.text_frame, "Leyes empíricas establecidas", 14, bold=True, color=ACCENT3)
    # Stefan-Boltzmann
    formula_img(s2, 'stefan', 0.6, 3.3, 2.4, 0.4)
    add_text(textbox(s2, 0.5, 3.25, 2.0, 0.3).text_frame, "Stefan-Boltzmann (1879):", 10, color=GRAY)
    formula_img(s2, 'sigma_val', 3.1, 3.35, 3.6, 0.3)
    # Wien
    add_text(textbox(s2, 0.5, 3.75, 2.0, 0.3).text_frame, "Wien (1893):", 10, color=GRAY)
    formula_img(s2, 'wien', 0.6, 3.95, 4.3, 0.35)

    # El desafío
    tb_d = textbox(s2, 0.4, 4.45, 6.0, 0.6)
    tf_d = tb_d.text_frame; tf_d.word_wrap = True
    add_text(tf_d, "El desafío", 14, bold=True, color=ACCENT3)
    add_text(tf_d, "La curva experimental era precisa y medible.", 12, color=WHITE)
    add_text(tf_d, "La energía total es finita — nunca infinita.", 12, bold=True, color=ACCENT2)

    # Analogía Teledetección
    rect(s2, 0.4, 5.2, 12.5, 0.03, ACCENT3)
    rect(s2, 0.4, 5.35, 12.5, 1.8, DARKBG)
    tb_a = textbox(s2, 0.55, 5.4, 12.2, 1.7)
    tf_a = tb_a.text_frame; tf_a.word_wrap = True
    add_text(tf_a, "🌍 Analogía en Teledetección / Radiometría", 13, bold=True, color=ACCENT3)
    add_text(tf_a, '"Un sensor multiespectral (Landsat / Sentinel-2) mide la radiancia espectral de la superficie.', 11, color=WHITE, italic=True)
    add_text(tf_a, 'Si la teoría clásica fuera cierta, cada banda espectral hacia el UV sumaría energía infinita.', 11, color=WHITE, italic=True)
    add_text(tf_a, 'Sin embargo, cualquier radiómetro real muestra que la señal decae en el UV — exactamente como predice Planck."', 11, color=WHITE, italic=True)

    # ═══════ SLIDE 3: CATÁSTROFE UV ═══════
    s3 = prs.slides.add_slide(blank)
    set_bg(s3)
    header(s3, "2. La Catástrofe del Ultravioleta: El Colapso Clásico", ACCENT2)

    s3.shapes.add_picture(CHART_RJ, Inches(6.4), Inches(1.15), Inches(6.6), Inches(3.8))

    tb3 = textbox(s3, 0.4, 1.15, 6.0, 1.0)
    tf3 = tb3.text_frame; tf3.word_wrap = True
    add_text(tf3, "La Ley de Rayleigh-Jeans (1900–05)", 14, bold=True, color=ACCENT2)
    add_text(tf3, "• Equipartición: cada modo recibe", 12, color=WHITE)

    # Formulas
    formula_img(s3, 'equip', 0.6, 1.8, 3.5, 0.35)
    add_text(textbox(s3, 0.4, 2.2, 6.0, 0.3).text_frame, "• Densidad de modos:", 12, color=WHITE)
    formula_img(s3, 'modes', 0.6, 2.45, 2.4, 0.5)

    add_text(textbox(s3, 0.4, 3.05, 6.0, 0.3).text_frame, "Ecuación Clásica de Rayleigh-Jeans:", 13, bold=True, color=ACCENT2)
    formula_img(s3, 'rj', 0.5, 3.35, 4.0, 0.5)

    # Crisis section
    tb_cr = textbox(s3, 0.4, 4.0, 6.0, 2.0)
    tf_cr = tb_cr.text_frame; tf_cr.word_wrap = True
    add_text(tf_cr, "La Crisis (Ehrenfest, 1911):", 13, bold=True, color=ACCENT2)
    add_text(tf_cr, "No es un error de cálculo.", 12, color=WHITE)
    add_text(tf_cr, "Es consecuencia inevitable de combinar", 12, color=WHITE)
    add_text(tf_cr, "termodinámica + electromagnetismo clásicos.", 12, color=WHITE)

    formula_img(s3, 'diverge', 0.6, 5.2, 3.6, 0.45)

    add_text(textbox(s3, 0.4, 5.75, 6.0, 0.5).text_frame,
             "∞ modos × energía constante = ∞", 14, bold=True, color=ACCENT2)

    # ═══════ SLIDE 4: SOLUCIÓN DE PLANCK ═══════
    s4 = prs.slides.add_slide(blank)
    set_bg(s4)
    header(s4, "3. La Hipótesis de Planck y la Cuantización de la Energía", ACCENT3)

    # Columna izquierda - Hipótesis
    tb_h = textbox(s4, 0.4, 1.15, 6.0, 0.5)
    add_text(tb_h.text_frame, "El Postulado de Cuantización (1900)", 14, bold=True, color=ACCENT3)
    add_text(tb_h.text_frame, "Los osciladores intercambian energía", 11, color=WHITE)
    add_text(tb_h.text_frame, "únicamente en paquetes discretos (cuantos):", 11, color=WHITE)

    formula_img(s4, 'quant', 0.5, 2.05, 5.0, 0.45)
    formula_img(s4, 'h_val', 0.6, 2.5, 3.2, 0.28)

    # Distribución
    add_text(textbox(s4, 0.4, 2.9, 6.0, 0.3).text_frame,
             "Distribución Cuántica de Planck:", 13, bold=True, color=ACCENT3)
    formula_img(s4, 'planck', 0.4, 3.2, 5.5, 0.55)

    # Supresión
    add_text(textbox(s4, 0.4, 3.9, 6.0, 0.3).text_frame,
             "Supresión Exponencial:", 13, bold=True, color=ACCENT3)
    formula_img(s4, 'low_freq', 0.5, 4.2, 4.8, 0.3)
    formula_img(s4, 'high_freq', 0.5, 4.55, 4.8, 0.3)

    # Gráfica
    s4.shapes.add_picture(CHART_RJ, Inches(6.4), Inches(1.15), Inches(6.6), Inches(3.2))

    # Cita + mecanismo
    tb_m = textbox(s4, 6.5, 4.45, 6.5, 0.9)
    tf_m = tb_m.text_frame; tf_m.word_wrap = True
    add_text(tf_m, 'El "precio de entrada" mínimo (hν) supera', 11, color=WHITE)
    add_text(tf_m, 'la energía térmica disponible (k_BT).', 11, color=WHITE)
    add_text(tf_m, '→ Los modos UV quedan "congelados".', 11, bold=True, color=ACCENT3)
    add_text(tf_m, '"Un acto de desesperación" — Max Planck', 10, italic=True, color=GRAY)

    # Analogía DN
    rect(s4, 0.4, 5.2, 12.5, 0.03, ACCENT3)
    rect(s4, 0.4, 5.35, 12.5, 1.8, DARKBG)
    tb_dn = textbox(s4, 0.55, 5.4, 12.2, 1.7)
    tf_dn = tb_dn.text_frame; tf_dn.word_wrap = True
    add_text(tf_dn, "📡 Analogía: Digitalización Radiométrica (DN)", 13, bold=True, color=ACCENT3)
    add_text(tf_dn, "En sensores de Teledetección (Sentinel-2, LiDAR), la radiancia analógica se muestrea en", 11, color=WHITE)
    add_text(tf_dn, "Valores Digitales Discretos (DN). No existe el DN 127.5, solo 127 o 128.", 11, color=WHITE)
    add_text(tf_dn, "Planck aplicó exactamente esta discretización a la energía: si k_BT no alcanza para", 11, color=WHITE)
    add_text(tf_dn, "subir al primer peldaño (hν), el modo queda desierto.", 11, bold=True, color=ACCENT3)

    # ═══════ SLIDE 5: CONCLUSIÓN E IMPACTO ═══════
    s5 = prs.slides.add_slide(blank)
    set_bg(s5)
    header(s5, "Conclusión: El Cambio de Paradigma en la Física", YELLOW)

    # Timeline
    timeline = [
        ("1879",    "Stefan",           "R = σT⁴",              ACCENT1),
        ("1893",    "Wien",             "λ_max·T = b",          ACCENT1),
        ("1900–05", "Rayleigh\n& Jeans", "Catástrofe UV",        ACCENT2),
        ("1900 ★",  "Planck",           "E = nhν",               ACCENT3),
        ("1905",    "Einstein",         "Fotones\nE = hν",       YELLOW),
        ("1913",    "Bohr",             "L = nℏ\nEₙ = −13.6/n²", ACCENT1),
        ("1925–26", "Heisenberg\nSchrödinger", "ĤΨ = EΨ",       PURPLE),
    ]
    x_start, x_step, y_line = 0.3, 1.8, 2.7
    rect(s5, x_start, y_line + 0.15, 12.5, 0.04, GRAY)

    for i, (yr, name, desc, col) in enumerate(timeline):
        x = x_start + i * x_step
        rect(s5, x, y_line + 0.03, 0.22, 0.28, col)
        tb_yr = textbox(s5, x-0.2, y_line-0.55, 1.5, 0.5)
        add_text(tb_yr.text_frame, yr, 12, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb_n = textbox(s5, x-0.3, y_line+0.45, 1.6, 0.55)
        add_text(tb_n.text_frame, name, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb_d = textbox(s5, x-0.35, y_line+1.0, 1.7, 0.7)
        add_text(tb_d.text_frame, desc, 9, color=GRAY, align=PP_ALIGN.CENTER)

    # Impacto tecnológico
    rect(s5, 0.4, 4.65, 12.5, 0.03, YELLOW)
    tb_imp = textbox(s5, 0.4, 4.8, 12.5, 0.4)
    add_text(tb_imp.text_frame, "⚡ Impacto en la Ciencia y Tecnología Moderna", 14, bold=True, color=YELLOW)

    impacts = [
        ("💡 Semiconductores",   "La banda prohibida es un\nefecto directo de la\ncuantización de la energía.", 0.4),
        ("📡 Sensores CMOS/CCD", "Cámaras satelitales (Sentinel,\nLandsat) detectan fotones\nindividuales.", 4.6),
        ("🔬 Láseres / LiDAR",  "LiDAR, Raman, fluorescencia:\ntodas las técnicas espectrales\ndependen de emisión cuántica.", 8.8),
    ]
    for label, desc, x in impacts:
        rect(s5, x, 5.3, 3.8, 1.15, DARKBG)
        tb_i = textbox(s5, x+0.15, 5.35, 3.5, 1.05)
        tf_i = tb_i.text_frame; tf_i.word_wrap = True
        add_text(tf_i, label, 12, bold=True, color=ACCENT3)
        add_text(tf_i, desc, 10, color=WHITE)

    # Frase de cierre
    rect(s5, 0.8, 6.55, 11.7, 0.04, YELLOW)
    tb_c = textbox(s5, 0.8, 6.65, 11.7, 0.6)
    tf_c = tb_c.text_frame; tf_c.word_wrap = True
    add_text(tf_c, '"La catástrofe ultravioleta no fue una falla menor — fue la grieta por donde entró toda la física cuántica."',
             14, italic=True, color=YELLOW, align=PP_ALIGN.CENTER)

    # Guardar
    prs.save(OUT_PPTX)
    sz = os.path.getsize(OUT_PPTX)
    print(f"\n✓ Presentación guardada en:\n  {OUT_PPTX}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size: {sz:,} bytes")

# ── MAIN ──
if __name__ == "__main__":
    print("Renderizando fórmulas LaTeX...")
    F = render_all_formulas()
    print("Generando gráficas HD...")
    make_blackbody_chart()
    make_rj_planck_chart()
    print("Construyendo presentación v4...")
    build_pptx()
